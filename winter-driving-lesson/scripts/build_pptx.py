#!/usr/bin/env python3
"""Build PowerPoint (.pptx) from the winter driving lesson."""

import html
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

BASE = Path(__file__).resolve().parent.parent
IMAGES = BASE / "images"
OUT = BASE / "נהיגת-חורף-בישראל.pptx"

# Colors matching the web lesson
TEAL = RGBColor(0x0D, 0x4F, 0x5C)
ORANGE = RGBColor(0xFF, 0x8C, 0x42)
MINT = RGBColor(0x2D, 0xD4, 0xBF)
DARK = RGBColor(0x0F, 0x24, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x9E, 0xC5, 0xDB)

LESSON = {
    "title": "נהיגת חורף בישראל",
    "subtitle": "מדריך לבטיחות בדרכים בגשם, ערפל וסופות — מותאם לישראל (ללא שלג)",
    "sources": [
        "משטרת ישראל — נהיגה בחורף",
        "משרד התחבורה והבטיחות בדרכים (gov.il)",
        "רלב\"ד Ask — נהיגת חורף",
    ],
    "chapters": [
        {
            "title": "חלק א׳ — הכנת הרכב לחורף",
            "tagline": "בדיקות לפני עונת הגשמים והערפל",
            "hero": "ch1-wipers.png",
            "sections": [
                {
                    "heading": "למה ההכנה חשובה בישראל?",
                    "image": "lesson-banner.png",
                    "bullets": [
                        "בישראל החורף: גשמים, ערפל, רוחות והצפות — לא שלג על הכביש",
                        "רכב לא מוכן = סיכון גבוה לתאונה (משטרת ישראל)",
                    ],
                },
                {
                    "heading": "צמיגים ולחץ אוויר",
                    "image": "ch1-tires.png",
                    "bullets": [
                        "עומק שקע מינימום 1.6 מ\"מ, ללא סדקים",
                        "לחץ אוויר לפי המלצת היצרן",
                        "צמיגים תקינים מקצרים מרחק בלימה על כביש רטוב",
                    ],
                },
                {
                    "heading": "מגבים, תאורה ומצבר",
                    "image": "ch1-wipers.png",
                    "bullets": [
                        "מגבים תקינים + מיכל נוזל שמשות מלא",
                        "כל הפנסים: אורות, בלימה, ערפל, מצוקה",
                        "מצבר — עומס חשמלי גבוה בחורף",
                        "בלמים והגה — תקינות לפני נסיעה",
                    ],
                },
                {
                    "heading": "שמשות, אדים ונראות",
                    "image": "ch1-defog.png",
                    "bullets": [
                        "ניקוי שמשות מבפנים ומבחוץ",
                        "מערכת הפשרה/מיזוג פועלת",
                        "חלון מעט פתוח — מפחית אדים",
                        "הסרת כובע/בגדים שמגבילים ראייה",
                    ],
                },
            ],
            "quiz": [
                (
                    "לפי משטרת ישראל — מה חשוב לבדוק בצמיגים?",
                    ["רק צבע הצמיג", "לחץ אוויר, שקע וסדקים", "רק גודל הרכב", "אין צורך בבדיקה"],
                    1,
                ),
                (
                    "תפקיד מגבים תקינים בחורף?",
                    ["מראה בלבד", "ראות טובה בגשם", "חיסכון בדלק", "רק בקיץ"],
                    1,
                ),
                (
                    "למה לבדוק תאורה לפני החורף?",
                    ["ראות מוגבלת בגשם/ערפל", "מחליף נהיגה זהירה", "רק בלילה", "אין חוק"],
                    0,
                ),
                (
                    "מה מפחית אדים על השמשה?",
                    ["חלונות סגורים לגמרי", "הפשרה + אוורור קל", "בגדים על השמשה", "כיבוי מגבים"],
                    1,
                ),
            ],
        },
        {
            "title": "חלק ב׳ — נהיגה בגשם",
            "tagline": "מהירות, מרחק, גשם ראשון והחלקה",
            "hero": "ch2-hero-rain.png",
            "sections": [
                {
                    "heading": "עקרונות הזהב",
                    "image": "ch2-distance.png",
                    "bullets": [
                        "האטה + מרחק — שני עקרונות מרכזיים (רלב\"ד)",
                        "גשם קל: מרחק מעקב ~100 מ'",
                        "גשם חזק: מרחק מעקב ~200 מ'",
                    ],
                },
                {
                    "heading": "גשם ראשון — הסכנה הנסתרת",
                    "image": "ch2-first-rain.png",
                    "bullets": [
                        "שמנים ואבק + מים ראשונים = כביש חלקלק במיוחד",
                        "להימנע מנהיגה מיותרת ברגעי הגשם הראשונים",
                        "אם נוסעים — מהירות נמוכה וזהירות יתרה",
                    ],
                },
                {
                    "heading": "אורות, בלימה והגה",
                    "image": "ch2-hero-rain.png",
                    "bullets": [
                        "הדלקת אורות גם ביום",
                        "שתי ידיים על ההגה",
                        "בלימה עדינה — לא פתאומית",
                        "הימנעות מעקיפות מסוכנות; זהירות מהולכי רגל",
                    ],
                },
                {
                    "heading": "שלוליות והידרופלנינג",
                    "image": "ch2-puddle.png",
                    "bullets": [
                        "לעקוף שלולית אם אפשר",
                        "אם חוצים: מהירות איטית, הגה ישר, תנועה רציפה",
                        "הידרופלנינג — האטה מראש",
                    ],
                },
            ],
            "quiz": [
                (
                    "שני עקרונות מפתח לנהיגה בגשם?",
                    ["מהירות ועקיפות", "האטה ומרחק", "כיבוי אורות", "נסיעה צמודה"],
                    1,
                ),
                (
                    "למה גשם ראשון מסוכן?",
                    ["כביש יבש", "שמנים ואבק יוצרים חלקלקות", "אין חוקים", "רק בלילה"],
                    1,
                ),
                (
                    "מתי להדליק אורות בגשם?",
                    ["רק בחניון", "גם ביום", "לעולם לא", "רק בעקיפה"],
                    1,
                ),
                (
                    "חציית שלולית ללא ברירה?",
                    ["מהירות גבוהה + פניה חדה", "איטי, הגה ישר, רציף", "עצירה באמצע", "לסגור עיניים"],
                    1,
                ),
            ],
        },
        {
            "title": "חלק ג׳ — ערפל, סופות והצפות",
            "tagline": "ראות מוגבלת, שיטפונים ונהיגה מול סיכון",
            "hero": "ch3-hero-fog.png",
            "sections": [
                {
                    "heading": "נהיגה בערפל",
                    "image": "ch3-hero-fog.png",
                    "bullets": [
                        "אור רגיל + פנסי ערפל — לא אור גבוה",
                        "האטה מאוד — ראייה קצרה",
                        "מרחק גדול; בלי עקיפות וצמוד מדי",
                        "ראות אפסית: עצירה בטוחה + מצוקה",
                    ],
                },
                {
                    "heading": "סופות, רוח ושיטפונים",
                    "image": "ch3-flood.png",
                    "bullets": [
                        "אין לחצות כבישים מוצפים, נחלים, שלוליות עמוקות",
                        "גם מעט מים עלול לסחוף רכב",
                        "זהירות מבורות מוסתרים מתחת למים",
                    ],
                },
                {
                    "heading": "מתי לעצור",
                    "image": "ch3-stop.png",
                    "bullets": [
                        "גשם/ערפל קיצוני — עצירה בצד, מצוקה, המתנה",
                        "הימנעות מעצירה על שולי כביש מסוכנים",
                    ],
                },
                {
                    "heading": "בטיחות כללית בחורף ישראלי",
                    "image": "lesson-banner.png",
                    "bullets": [
                        "אין נהיגת שלג — כן: גשם, ערפל, רוח, הצפות",
                        "קור רוח, לא לנהוג בעייפות",
                        "מידע רשמי: gov.il, משטרת ישראל",
                    ],
                },
            ],
            "quiz": [
                (
                    "איזה אור בערפל?",
                    ["אור גבוה בלבד", "רגיל + ערפל, לא גבוה", "ללא אורות", "רק חניה"],
                    1,
                ),
                (
                    "כבישים מוצפים — משטרת ישראל?",
                    ["מעודדת מרוץ", "לא לחצות מים ונחלים", "מבטלת חוקים", "נסיעה מהירה"],
                    1,
                ),
                (
                    "ראות אפסית בערפל?",
                    ["מהירות מלאה", "עצירה בטוחה + מצוקה", "כיבוי אורות", "עקיפה בלי מראות"],
                    1,
                ),
                (
                    "לא רלוונטי לחורף בישראל?",
                    ["גשם", "ערפל והצפות", "נהיגה על שלג וקרח", "מרחק בראות מוגבלת"],
                    2,
                ),
            ],
        },
    ],
}


def strip_html(text: str) -> str:
    text = re.sub(r"<[^>]+>", "", text)
    return html.unescape(text).strip()


def set_rtl_paragraph(paragraph, align=PP_ALIGN.RIGHT):
    paragraph.alignment = align
    try:
        p_pr = paragraph._p.get_or_add_pPr()
        p_pr.set("rtl", "1")
    except Exception:
        pass


def add_title_slide(prs: Presentation):
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = LESSON["title"]
    slide.placeholders[1].text = LESSON["subtitle"]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for p in shape.text_frame.paragraphs:
            set_rtl_paragraph(p)
            for run in p.runs:
                run.font.name = "Arial"
                run.font.size = Pt(28 if shape == slide.shapes.title else 18)
    img = IMAGES / "lesson-banner.png"
    if img.exists():
        slide.shapes.add_picture(str(img), Inches(0.5), Inches(3.8), width=Inches(9))


def add_section_slide(prs: Presentation, title: str, subtitle: str = "", image_name: str | None = None):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = TEAL
    bg.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    set_rtl_paragraph(p)
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.font.name = "Arial"

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        set_rtl_paragraph(p2)
        p2.font.size = Pt(20)
        p2.font.color.rgb = MINT
        p2.font.name = "Arial"

    if image_name:
        path = IMAGES / image_name
        if path.exists():
            slide.shapes.add_picture(str(path), Inches(0.5), Inches(1.8), width=Inches(9))
    return slide


def add_content_slide(prs: Presentation, heading: str, bullets: list[str], image_name: str | None = None):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    accent = slide.shapes.add_shape(1, 0, 0, Inches(0.15), prs.slide_height)
    accent.fill.solid()
    accent.fill.fore_color.rgb = ORANGE
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.7))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = heading
    set_rtl_paragraph(tp)
    tp.font.bold = True
    tp.font.size = Pt(26)
    tp.font.color.rgb = TEAL
    tp.font.name = "Arial"

    text_left = Inches(0.4)
    text_width = Inches(4.8)
    img_left = Inches(5.4)
    img_width = Inches(4.1)

    if image_name and (IMAGES / image_name).exists():
        slide.shapes.add_picture(str(IMAGES / image_name), img_left, Inches(1.1), width=img_width)
    else:
        text_width = Inches(9.2)

    body = slide.shapes.add_textbox(text_left, Inches(1.0), text_width, Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        set_rtl_paragraph(p)
        p.font.size = Pt(17)
        p.font.color.rgb = DARK
        p.font.name = "Arial"
        p.space_after = Pt(8)


def add_quiz_slide(prs: Presentation, part: str, num: int, question: str, options: list[str], correct: int):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    header = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.5))
    hp = header.text_frame.paragraphs[0]
    hp.text = f"{part} — מבחן שאלה {num}/4"
    set_rtl_paragraph(hp)
    hp.font.size = Pt(14)
    hp.font.color.rgb = ORANGE
    hp.font.bold = True

    qbox = slide.shapes.add_textbox(Inches(0.4), Inches(0.75), Inches(9.2), Inches(1.2))
    qp = qbox.text_frame.paragraphs[0]
    qp.text = question
    set_rtl_paragraph(qp)
    qp.font.size = Pt(22)
    qp.font.bold = True
    qp.font.color.rgb = TEAL

    letters = ["א", "ב", "ג", "ד"]
    obox = slide.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(8.8), Inches(4.5))
    tf = obox.text_frame
    for i, opt in enumerate(options):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        mark = "✓ " if i == correct else ""
        p.text = f"{mark}{letters[i]}. {opt}"
        set_rtl_paragraph(p)
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0x14, 0xB8, 0xA6) if i == correct else DARK

    note = slide.shapes.add_textbox(Inches(0.4), Inches(6.5), Inches(9.2), Inches(0.4))
    np = note.text_frame.paragraphs[0]
    np.text = "תשובה נכונה מסומנת ב-✓ (באתר: תשובה שגויה = חזרה על החלק)"
    set_rtl_paragraph(np)
    np.font.size = Pt(11)
    np.font.color.rgb = MUTED


def add_sources_slide(prs: Presentation):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "מקורות רשמיים"
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, src in enumerate(LESSON["sources"]):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = f"• {src}"
        set_rtl_paragraph(p)
        p.font.size = Pt(20)
    p = body.add_paragraph()
    p.text = "🇮🇱 מותאם לישראל — ללא נהיגת שלג"
    set_rtl_paragraph(p)
    p.font.size = Pt(18)
    p.font.color.rgb = ORANGE


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)

    add_section_slide(
        prs,
        "מבנה השיעור",
        "3 חלקים + מבחן של 4 שאלות בסוף כל חלק",
        "lesson-banner.png",
    )

    overview = slide = prs.slides.add_slide(prs.slide_layouts[1])
    overview.shapes.title.text = "תוכן העניינים"
    tf = overview.placeholders[1].text_frame
    for i, ch in enumerate(LESSON["chapters"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{i + 1}. {ch['title']} — {ch['tagline']}"
        set_rtl_paragraph(p)
        p.font.size = Pt(20)

    for ch in LESSON["chapters"]:
        add_section_slide(prs, ch["title"], ch["tagline"], ch.get("hero"))
        for sec in ch["sections"]:
            add_content_slide(prs, sec["heading"], sec["bullets"], sec.get("image"))
        add_section_slide(prs, f"מבחן — {ch['title']}", "4 שאלות רב־ברירה (א–ד)")
        for qi, (q, opts, correct) in enumerate(ch["quiz"], 1):
            add_quiz_slide(prs, ch["title"], qi, q, opts, correct)

    add_sources_slide(prs)

    # Closing
    add_section_slide(prs, "סיום", "נהיגה בטוחה וחורף שקט! 🌧️", "lesson-banner.png")

    prs.save(OUT)
    print(f"Saved: {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
